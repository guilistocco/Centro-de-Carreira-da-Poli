arquivo_excel = r'CV e Foto para Pitch Pessoal - Programa de Carreira 2020 (respostas).xlsx' #ok
arquivo_ppt_template = r'CC_Template_PPT_20180724 - Copia.pptx' #ok
titulo_programa = "Programa de Carreira 2020" #ok

#bibliotecas de python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.text.text import Font as font

# pandas para organização e carregamento dos dados
import pandas as pd
import numpy as np

# OS para abrirmos a apresentação assim que terminar
import os

#
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive

#quando colocamos um argumento, ele busca por esse pptx no diretorio
prs = Presentation('CC_Template_PPT_20180724 - Copia.pptx')


# cria um slide com o layout 0 das opções, nesse template é o com titulo e subtítulo
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

# coloca o titulo e subtitulo que queremos
title.text = "Portfólio de Alunos"
subtitle.text = titulo_programa

# Directory 
directory_fotos = "Fotos formulario"
directory_CV = "Resume Book"
  
# Diretorio atual
parent_dir = os.getcwd()

# Path 
caminho_fotos = os.path.join(parent_dir, directory_fotos) 
caminho_CV = os.path.join(parent_dir, directory_CV) 

if  os.path.isdir(caminho_fotos) == False:
    os.mkdir(caminho_fotos) 

if  os.path.isdir(caminho_CV) == False:
    os.mkdir(caminho_CV) 

gauth = GoogleAuth()
gauth.LocalWebserverAuth()
drive = GoogleDrive(gauth)


f1 = pd.read_excel(arquivo_excel,
                     header=0,
                     parse_dates=True)
f1 = f1.sort_values('Nome completo').reset_index()


# cria uma coluna com somente o ID do arquivo uploadado para futuro download
fotos_id = list(f1['Foto'].values)

f_ids = []

for f in fotos_id:
    f_ids.append(f[33:])
f1['foto_id'] = f_ids

##########################################################
CV_id = list(f1['CV'].values)

CV_ids = []

for f in CV_id:
    CV_ids.append(f[33:])
f1['foto_id'] = CV_ids



# percorrendo todas as linhas do df

for row in f1.index:
    
    # cria um slide com o layout 2 das opções, nesse template é o com titulo e texto
    texto_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(texto_slide_layout)

    #coloca o nome como título
    title = slide.shapes.title
    title.text = str(f1["Nome completo"][row])

    # cria caixa de texto add_textbox(left, top, width, height)
    txBox = slide.shapes.add_textbox(Cm(15),Cm(5),Cm(11),Cm(11))
    tf = txBox.text_frame

    #coloca email e outras informações
    #------------------- EMAIL -----------------
    p = tf.add_paragraph()
    p = tf.paragraphs[1]
    run = p.add_run()
    font = run.font
    font.bold = True
    run.text = "Email: " 
    run = p.add_run()
    run.text = str(f1['E-mail'][row])

    #------------------- TELEFONE -----------------

    p = tf.add_paragraph()
    p = tf.paragraphs[2]
    run = p.add_run()
    font = run.font
    font.bold = True
    run.text = "Telefone: " 
    run = p.add_run()
    run.text = str(f1['Telefone'][row])

    #------------------- CURSO -----------------


    p = tf.add_paragraph()
    p = tf.paragraphs[3]
    run = p.add_run()
    font = run.font
    font.bold = True
    run.text = "Curso: "
    run = p.add_run()
    run.text = str(f1['Curso'][row])
    
    #------------------- FORMATURA -----------------

    p = tf.add_paragraph()
    p = tf.paragraphs[4]
    run = p.add_run()
    font = run.font
    font.bold = True
    run.text = "Ano de Formatura: "
    run = p.add_run()
    run.text = str(f1['Ano de Formatura'][row])
    


    #------------------- ESTÁGIO -----------------


    p = tf.add_paragraph()
    p = tf.paragraphs[5]
    run = p.add_run()
    font = run.font
    font.bold = True
    run.text = "Disponibilidade para Estágio Regular: "
    run = p.add_run()
    run.text = str(f1['Disponibilidade para Estágio Regular ao fim do Programa'][row])
        
    p = tf.add_paragraph()
    p = tf.paragraphs[6]
    run = p.add_run()
    font = run.font
    font.bold = True
    run.text = "Disponibilidade para Estágio de Férias: "
    run = p.add_run()
    run.text = str(f1['Disponibilidade para Estágio de Férias ao fim do Programa'][row])

    p = tf.add_paragraph()
    p = tf.paragraphs[7]
    run = p.add_run()
    font = run.font
    font.bold = True
    run.text = "Disponibilidade para Estágio Quadrimestral: "
    run = p.add_run()
    run.text = str(f1['Disponibilidade para Estágio Quadrimestral ao fim do Programa'][row])

    #------------------- LINKEDIN -----------------

    p = tf.add_paragraph()
    p = tf.paragraphs[8]
    run = p.add_run()
    font = run.font
    font.bold = True
    run.text = "LinkedIn: "

    run = p.add_run()
    font = run.font
    font.size = Pt(12)
    run.hyperlink.address = f1['Link para o Linkedin'][row]
    run.text = f1['Link para o Linkedin'][row]

       
    #------------------- CURRICULO -----------------

    CV_id = CV_ids[row]
    destination_CV = caminho_CV + '//' + str(f1["Nome completo"][row]) + '_CV.pdf'
    CV_path = directory_CV +'//' + str(f1["Nome completo"][row]) + '_CV.pdf'
    file6 = drive.CreateFile({'id': CV_id})
    file6.GetContentFile(CV_path)

#    p = tf.add_paragraph()
#    p = tf.paragraphs[9]
#   run = p.add_run()
#    font = run.font
#    font.bold = True
#    run.text = "Link para Currículo: "
    
#    run = p.add_run()
#    font = run.font
#    font.size = Pt(10)
#    run.hyperlink.address = str(f1['CV'][row])
#    run.text = str(f1['CV'][row])


    ## download da foto da pessoa
    #------------------- FOTO -----------------

    file_id = f_ids[row]
    destination = caminho_fotos + '//' + str(f1["Nome completo"][row]) + '_foto.jpg'
    img_path = directory_fotos +'//' + str(f1["Nome completo"][row]) + '_foto.jpg'
    file6 = drive.CreateFile({'id': file_id})
    file6.GetContentFile(img_path)
    
    try:
        #add_picture(img_path, left, top, height=height)
        pic = slide.shapes.add_picture(img_path, Cm(3) , Cm(5), height=Cm(12), width=Cm(12))
    except:
        pass

prs.save('final.pptx')
os.startfile('final.pptx')