# -*- coding: utf-8 -*-
"""
Created on Wed Jul 15 22:21:17 2020

@author: Windows10
"""


#importando as bibliotecas nessários

#bibliotecas de python-pptx
from pptx import Presentation
from pptx.util import Cm

# pandas para organização e carregamento dos dados
import pandas as pd

# OS para abrirmos a apresentação assim que terminar
import os

# requests nos ajuda a buscar as fotos no google drive
import requests

#funções fazem o liberam as autorizações e baixam o conteúdo
def download_file_from_google_drive(id, destination):
    URL = "https://docs.google.com/uc?export=download"

    session = requests.Session()

    response = session.get(URL, params = { 'id' : id }, stream = True)
    token = get_confirm_token(response)

    if token:
        params = { 'id' : id, 'confirm' : token }
        response = session.get(URL, params = params, stream = True)

    save_response_content(response, destination)    

def get_confirm_token(response):
    for key, value in response.cookies.items():
        if key.startswith('download_warning'):
            return value

    return None

def save_response_content(response, destination):
    CHUNK_SIZE = 32768

    with open(destination, "wb") as f:
        for chunk in response.iter_content(CHUNK_SIZE):
            if chunk: # filter out keep-alive new chunks
                f.write(chunk)
                

parent_dir = os.getcwd()
        

#quando colocamos um argumento, ele busca por esse pptx no diretorio
prs = Presentation(parent_dir + '\\' + 'CC_Template_PPT_20180724 - Copia.pptx')

    # cria um slide com o layout 0 das opções, nesse template é o com titulo e subtítulo
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

# coloca o titulo e subtitulo que queremos
title.text = "Carômetro"
subtitle.text = "Programa de ________ 20XX"

# Python program to explain os.mkdir() method 
  
# Directory 
directory = "Fotos_formulario"
  
# Diretorio atual

# Path 
caminho = os.path.join(parent_dir, directory) 
 
if  os.path.isdir(caminho) == False:
    os.mkdir(caminho) 

# na frente do r, coloca o caminho da planilha do formulário
forms = pd.read_excel(parent_dir + '\\' + 'Carometro_TESTE_(respostas).xlsx',
                     header=0,
                     usecols=['Nome','Data_de_nascimento','Engenharia','NUSP','Endereço de e-mail','Foto para o Carômetro'],
                     parse_dates=True) #cuidado com essas colunas, podem mudar de uma planilha pra outra

# coloca a data de nascimento no formato DD/MM/AAAA
forms['Data_de_nascimento'] = forms['Data_de_nascimento'].dt.strftime('%d/%m/%Y')


# cria uma coluna com somente o ID do arquivo uploadado para futuro download
fotos_id = list(forms['Foto para o Carômetro'].values)
f_ids = []
for f in fotos_id:
    f_ids.append(f[33:])
forms['foto_id'] = f_ids
    

# percorrendo todas as linhas do df
for row in forms.index:
    
    # cria um slide com o layout 2 das opções, nesse template é o com titulo e texto
    texto_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(texto_slide_layout)

    #coloca o nome como título
    title = slide.shapes.title
    title.text = str(forms.Nome[row])

    # cria caixa de texto
    txBox = slide.shapes.add_textbox(Cm(19),Cm(8),Cm(11),Cm(4))
    tf = txBox.text_frame

    #coloca email e outras informações
    p = tf.add_paragraph()
    p.text = "Email: " + str(forms['Endereço de e-mail'][row])

    p = tf.add_paragraph()
    p.text = "Data de Nascimento: " + str(forms.Data_de_nascimento[row])

    p = tf.add_paragraph()
    p.text = "Engenharia " + str(forms.Engenharia[row])

    ## download da foto da pessoa, usando as funções já declaradas
    if __name__ == "__main__":
        file_id = f_ids[row]
        destination = caminho + '//' + str(forms.Nome[row]) + '_foto.jpg'
        download_file_from_google_drive(file_id, destination)

    img_path = caminho + '//' + str(forms.Nome[row]) + '_foto.jpg'
    
    #adiciona a foto na apresentação
    pic = slide.shapes.add_picture(img_path, Cm(6) , Cm(5), height=Cm(12), width=Cm(12))

    
# salva e abre a apresentação
prs.save('teste_6.pptx')
os.startfile('teste_6.pptx')